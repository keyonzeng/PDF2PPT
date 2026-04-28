import os
import subprocess
import tempfile
import time
from pathlib import Path
from uuid import uuid4

from PIL import Image, ImageChops, ImageFilter, ImageStat
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.mineru_service import process_pdf
from app.services.parser_service import parse_mineru_output
from app.services.ppt_gen_service import generate_pptx

try:
    import pypdfium2 as pdfium
except ImportError:  # pragma: no cover - optional runtime dependency
    pdfium = None

BASE_DIR = Path(__file__).resolve().parents[1]
SAMPLE_PDF = BASE_DIR / "uploads" / "TheLastLeaf.pdf"
OPENCLAW_PDF = BASE_DIR / "uploads" / "openclaw.pdf"
ACTION_SYSTEM_PDF = BASE_DIR / "uploads" / "action_system.pdf"


def _export_pptx_to_pdf(pptx_path: Path, output_pdf: Path, timeout_seconds: int = 180) -> None:
    script = """
$pptxPath = $env:PDF2PPT_PPTX_PATH
$pdfPath = $env:PDF2PPT_PDF_PATH

$powerPoint = $null
$presentation = $null
try {
    $powerPoint = New-Object -ComObject PowerPoint.Application
    $presentation = $powerPoint.Presentations.Open($pptxPath, $false, $false, $false)
    $presentation.SaveAs($pdfPath, 32)
}
finally {
    if ($presentation -ne $null) {
        $presentation.Close() | Out-Null
    }
    if ($powerPoint -ne $null) {
        $powerPoint.Quit() | Out-Null
    }
}
"""

    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile("w", suffix=".ps1", delete=False, encoding="utf-8") as handle:
        script_path = Path(handle.name)
        handle.write(script)

    last_result = None
    try:
        for attempt in range(3):
            try:
                result = subprocess.run(
                    [
                        "powershell",
                        "-NoProfile",
                        "-ExecutionPolicy",
                        "Bypass",
                        "-File",
                        str(script_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=timeout_seconds,
                    check=False,
                    env={**os.environ, "PDF2PPT_PPTX_PATH": str(pptx_path), "PDF2PPT_PDF_PATH": str(output_pdf)},
                )
            except subprocess.TimeoutExpired as exc:
                last_result = exc
                time.sleep(2)
                continue

            last_result = result
            if result.returncode == 0 and output_pdf.exists():
                return
            if attempt < 2:
                time.sleep(2)
                continue
            break
    finally:
        try:
            script_path.unlink(missing_ok=True)
        except OSError:
            pass

    if isinstance(last_result, subprocess.TimeoutExpired):
        raise RuntimeError(f"PowerPoint export timed out after retries: {last_result}")
    if last_result is not None and getattr(last_result, "returncode", 1) != 0:
        raise RuntimeError(
            "PowerPoint export failed: "
            f"returncode={last_result.returncode} stdout={last_result.stdout!r} stderr={last_result.stderr!r}"
        )
    raise RuntimeError(
        "PowerPoint export did not create PDF after retries: "
        f"stdout={getattr(last_result, 'stdout', '')!r} stderr={getattr(last_result, 'stderr', '')!r}"
    )


def _render_pdf_page(pdf_path: Path, page_index: int, scale: float = 2.0) -> Image.Image:
    if pdfium is None:
        raise RuntimeError("pypdfium2 is required for visual regression tests")
    document = pdfium.PdfDocument(str(pdf_path))
    try:
        page = document[page_index]
        bitmap = page.render(scale=scale)
        return bitmap.to_pil()
    finally:
        close = getattr(document, "close", None)
        if callable(close):
            close()


def _content_bbox(image: Image.Image, white_threshold: int = 245, padding: int = 12) -> tuple[int, int, int, int] | None:
    grayscale = image.convert("L")
    inverted = ImageChops.invert(grayscale)
    bbox = inverted.point(lambda value: 255 if value > (255 - white_threshold) else 0).getbbox()
    if bbox is None:
        return None
    left, top, right, bottom = bbox
    left = max(left - padding, 0)
    top = max(top - padding, 0)
    right = min(right + padding, image.width)
    bottom = min(bottom + padding, image.height)
    if right <= left or bottom <= top:
        return None
    return left, top, right, bottom


def _page_similarity(expected_image: Image.Image, actual_image: Image.Image) -> float:
    target_size = (192, 108)
    expected = expected_image.convert("L").resize(target_size, Image.Resampling.LANCZOS).filter(ImageFilter.GaussianBlur(radius=2.6))
    actual = actual_image.convert("L").resize(target_size, Image.Resampling.LANCZOS).filter(ImageFilter.GaussianBlur(radius=2.6))
    expected_bbox = _content_bbox(expected)
    actual_bbox = _content_bbox(actual)
    if expected_bbox or actual_bbox:
        if expected_bbox and actual_bbox:
            left = min(expected_bbox[0], actual_bbox[0])
            top = min(expected_bbox[1], actual_bbox[1])
            right = max(expected_bbox[2], actual_bbox[2])
            bottom = max(expected_bbox[3], actual_bbox[3])
        else:
            left, top, right, bottom = expected_bbox or actual_bbox  # type: ignore[assignment]
        expected = expected.crop((left, top, right, bottom))
        actual = actual.crop((left, top, right, bottom))
    diff = ImageChops.difference(expected, actual)
    stats = ImageStat.Stat(diff)
    mean_diff = stats.mean[0] if stats.mean else 0.0
    return max(0.0, 1.0 - (mean_diff / 550.0))


def _assert_pdf_visual_similarity(source_pdf: Path, generated_pptx: Path, sample_name: str, threshold: float = 0.90) -> None:
    with tempfile.TemporaryDirectory(prefix=f"visual-{sample_name}-") as tmp_dir:
        rendered_pdf = Path(tmp_dir) / f"{sample_name}.pdf"
        _export_pptx_to_pdf(generated_pptx, rendered_pdf)

        source_document = pdfium.PdfDocument(str(source_pdf)) if pdfium is not None else None
        generated_document = pdfium.PdfDocument(str(rendered_pdf)) if pdfium is not None else None
        assert source_document is not None and generated_document is not None

        try:
            assert len(source_document) == len(generated_document), (
                f"{sample_name} page count mismatch: source={len(source_document)} generated={len(generated_document)}"
            )

            page_scores: list[float] = []
            for page_index in range(len(source_document)):
                source_image = _render_pdf_page(source_pdf, page_index)
                generated_image = _render_pdf_page(rendered_pdf, page_index)
                similarity = _page_similarity(source_image, generated_image)
                page_scores.append(similarity)
                assert similarity >= threshold, (
                    f"{sample_name} page {page_index + 1} visual similarity below target: {similarity:.3f}"
                )

            print(f"visual_page_scores[{sample_name}]", [round(score, 3) for score in page_scores])
        finally:
            for document in (source_document, generated_document):
                close = getattr(document, "close", None)
                if callable(close):
                    close()


def _generated_content_slides(ppt: PptxPresentation, parsed_slide_count: int):
    assert len(ppt.slides) >= parsed_slide_count, "Generated PPT should include all parsed slides"
    generated_slide_delta = len(ppt.slides) - parsed_slide_count
    assert generated_slide_delta <= 1, "Generated PPT should not introduce more than one template slide"
    return list(ppt.slides)[-parsed_slide_count:]


def _normalize_shape_bbox(shape, slide_width: int, slide_height: int) -> tuple[float, float, float, float]:
    return (
        shape.left / slide_width,
        shape.top / slide_height,
        shape.width / slide_width,
        shape.height / slide_height,
    )


def _normalize_element_bbox(element, slide_width: float, slide_height: float) -> tuple[float, float, float, float]:
    x0, y0, x1, y1 = element.bbox
    return (
        x0 / slide_width,
        y0 / slide_height,
        (x1 - x0) / slide_width,
        (y1 - y0) / slide_height,
    )


def _assert_bbox_close(actual: tuple[float, float, float, float], expected: tuple[float, float, float, float], tolerance: float = 0.03):
    for actual_value, expected_value in zip(actual, expected):
        assert abs(actual_value - expected_value) <= tolerance, (
            f"Bounding box mismatch. actual={actual}, expected={expected}, tolerance={tolerance}"
        )


def _bbox_distance(actual: tuple[float, float, float, float], expected: tuple[float, float, float, float]) -> float:
    return sum(abs(actual_value - expected_value) for actual_value, expected_value in zip(actual, expected))


def _normalized_text(value: str) -> str:
    return " ".join(value.split())


def _find_text_shape(slide, expected_text: str, expected_bbox: tuple[float, float, float, float], used_shape_ids: set[int], slide_width: int, slide_height: int):
    normalized_expected = " ".join(expected_text.split())
    matched_candidates = []
    for shape in slide.shapes:
        if id(shape) in used_shape_ids:
            continue
        if hasattr(shape, "text") and shape.text:
            normalized_actual = " ".join(shape.text.split())
            if normalized_actual == normalized_expected:
                actual_bbox = _normalize_shape_bbox(shape, slide_width, slide_height)
                matched_candidates.append((shape, _bbox_distance(actual_bbox, expected_bbox)))

    if not matched_candidates:
        return None

    matched_candidates.sort(key=lambda item: item[1])
    return matched_candidates[0][0]


def _picture_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _has_full_page_snapshot(slide, slide_width: int, slide_height: int, threshold: float = 0.85) -> bool:
    for shape in _picture_shapes(slide):
        if shape.left == 0 and shape.top == 0 and shape.width == slide_width and shape.height == slide_height:
            return True
    return False


def _run_pipeline(sample_pdf: Path, request_prefix: str):
    assert sample_pdf.exists(), f"Sample PDF not found: {sample_pdf}"

    process_result = process_pdf(str(sample_pdf), request_id=request_prefix)
    assert process_result["status"] == "success", process_result

    output_folder = Path(process_result["output_folder"])
    assert output_folder.exists(), f"MinerU output folder missing: {output_folder}"

    content_list_candidates = list(output_folder.rglob("*_content_list.json"))
    assert content_list_candidates, f"No MinerU content list found under: {output_folder}"

    presentation = parse_mineru_output(str(output_folder))
    request_id = f"{request_prefix}-{uuid4().hex[:8]}"
    pptx_path = Path(
        generate_pptx(
            presentation,
            template_key="default",
            request_id=request_id,
            source_pdf_path=str(sample_pdf),
        )
    )
    assert pptx_path.exists(), f"Generated PPTX missing: {pptx_path}"

    ppt = PptxPresentation(str(pptx_path))
    content_slides = _generated_content_slides(ppt, len(presentation.slides))
    return presentation, ppt, content_slides, output_folder, pptx_path


def test_real_pdf_to_ppt_pipeline():
    presentation, ppt, _, output_folder, pptx_path = _run_pipeline(SAMPLE_PDF, "real-pipeline-test")

    assert output_folder.exists(), f"MinerU output folder missing: {output_folder}"
    assert pptx_path.exists(), f"Generated PPTX missing: {pptx_path}"
    assert len(presentation.slides) >= 2, "Expected at least 2 slides from sample PDF"
    assert any(element.type == "text" for slide in presentation.slides for element in slide.elements), "Expected parsed text elements"

    all_text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text.append(shape.text)

    combined_text = "\n".join(all_text)
    # OCR may produce artificially spaced characters; after cleanup the title can become "TheLastLeaf"
    assert ("The Last Leaf" in combined_text or "TheLastLeaf" in combined_text), "Expected representative title text in generated PPT"


def test_openclaw_pdf_content_and_layout_consistency():
    presentation, ppt, content_slides, output_folder, pptx_path = _run_pipeline(OPENCLAW_PDF, "openclaw-pipeline-test")

    assert output_folder.exists(), f"MinerU output folder missing: {output_folder}"
    assert pptx_path.exists(), f"Generated PPTX missing: {pptx_path}"
    assert presentation.slides, "Expected parsed slides from openclaw.pdf"

    slide_width = ppt.slide_width
    slide_height = ppt.slide_height

    matched_text_count = 0
    matched_picture_count = 0
    total_text_count = 0
    total_picture_like_count = 0

    for parsed_slide, generated_slide in zip(presentation.slides, content_slides):
        expected_text_elements = [
            element for element in parsed_slide.elements if element.type == "text" and element.content.strip() and element.bbox
        ]
        expected_picture_elements = [
            element for element in parsed_slide.elements
            if element.type in {"image", "table"} and element.bbox
        ]

        total_text_count += len(expected_text_elements)
        total_picture_like_count += len(expected_picture_elements)

        generated_pictures = _picture_shapes(generated_slide)
        has_full_page_snapshot = _has_full_page_snapshot(generated_slide, slide_width, slide_height)
        assert not has_full_page_snapshot, "Generated slide should not rely on a full-page snapshot picture"
        assert len(generated_pictures) >= len(expected_picture_elements), (
            "Generated slide should contain at least the parsed image/table elements"
        )

        unused_pictures = generated_pictures.copy()
        generated_text_entries = []
        for shape in generated_slide.shapes:
            if hasattr(shape, "text") and shape.text:
                generated_text_entries.append(
                    {
                        "shape": shape,
                        "text": _normalized_text(shape.text),
                        "bbox": _normalize_shape_bbox(shape, slide_width, slide_height),
                    }
                )
        used_text_entry_indexes = set()

        for element in expected_text_elements:
            expected_bbox = _normalize_element_bbox(element, parsed_slide.width, parsed_slide.height)
            normalized_expected_text = _normalized_text(element.content)
            text_candidates = []
            for entry_index, entry in enumerate(generated_text_entries):
                if entry_index in used_text_entry_indexes:
                    continue
                if entry["text"] == normalized_expected_text:
                    text_candidates.append((entry_index, entry, _bbox_distance(entry["bbox"], expected_bbox)))

            assert text_candidates, f"Missing text element in generated PPT: {element.content[:80]}"
            text_candidates.sort(key=lambda item: item[2])
            matched_entry_index, matched_entry, _ = text_candidates[0]

            _assert_bbox_close(matched_entry["bbox"], expected_bbox)
            used_text_entry_indexes.add(matched_entry_index)
            matched_text_count += 1

        for element in expected_picture_elements:
            expected_bbox = _normalize_element_bbox(element, parsed_slide.width, parsed_slide.height)
            matched_picture = None
            for picture in unused_pictures:
                actual_bbox = _normalize_shape_bbox(picture, slide_width, slide_height)
                if all(abs(actual_value - expected_value) <= 0.03 for actual_value, expected_value in zip(actual_bbox, expected_bbox)):
                    matched_picture = picture
                    break

            assert matched_picture is not None, "Missing image/table element with matching layout in generated PPT"
            unused_pictures.remove(matched_picture)
            matched_picture_count += 1

    assert total_text_count > 0, "Expected parsed text elements from openclaw.pdf"
    assert matched_text_count == total_text_count, "All parsed text elements should be preserved in generated PPT"
    assert matched_picture_count == total_picture_like_count, "All parsed image/table elements should keep layout mapping in generated PPT"


def test_real_pdf_to_ppt_visual_similarity():
    for sample_pdf, request_prefix in [
        (SAMPLE_PDF, "visual-the-last-leaf"),
        (ACTION_SYSTEM_PDF, "visual-action-system"),
    ]:
        presentation, _, _, _, pptx_path = _run_pipeline(sample_pdf, request_prefix)
        assert len(presentation.slides) > 0, f"Expected parsed slides from {sample_pdf.name}"
        _assert_pdf_visual_similarity(sample_pdf, pptx_path, sample_pdf.stem)
