from app.services.ppt_gen_service import generate_pptx, _element_geometry, _build_text_masked_picture_stream
from app.core.models import Presentation, Slide, TextElement, DocumentStyleProfile
from app.core.models import ImageElement
from app.main import _presentation_summary
import os
from pathlib import Path
from pptx import Presentation as PptxPresentation
from PIL import Image, ImageDraw

def test_generate_pptx_simple():
    # Arrange
    presentation = Presentation()
    slide1 = Slide(page_id=1, width=595.0, height=842.0)
    slide1.add_element(TextElement(content="Test Title", bbox=[50.0, 50.0, 500.0, 100.0], type="text"))
    slide1.add_element(TextElement(content="Body Content", bbox=[50.0, 150.0, 500.0, 300.0], type="text"))
    presentation.slides.append(slide1)
    
    # Act
    output_path = generate_pptx(presentation, template_key="default", request_id="test-run")
    
    # Assert
    assert os.path.exists(output_path)
    assert output_path.endswith(".pptx")
    print(f"Generated PPTX at: {output_path}")


def test_presentation_summary_keeps_auto_for_low_complexity_visual_slide():
    presentation = Presentation()
    slide = Slide(page_id=1, width=1376.0, height=768.0)
    slide.add_element(TextElement(content="Title", bbox=[80.0, 40.0, 500.0, 120.0], semantic_role="title"))
    slide.add_element(TextElement(content="Short body", bbox=[80.0, 160.0, 500.0, 240.0], semantic_role="body"))
    slide.add_element(ImageElement(content="[IMAGE]", path=__file__, bbox=[700.0, 120.0, 1200.0, 620.0]))
    presentation.slides.append(slide)

    summary = _presentation_summary(presentation, "req-auto", "sample.pdf", "output-folder")

    assert summary["slides"][0]["default_render_mode"] == "auto"


def test_presentation_summary_uses_image_fallback_for_fragmented_high_complexity_slide():
    presentation = Presentation()
    slide = Slide(page_id=1, width=1376.0, height=768.0)
    slide.add_element(ImageElement(content="[IMAGE]", path=__file__, bbox=[0.0, 0.0, 1376.0, 768.0]))
    for index in range(8):
        top = 40.0 + (index * 70.0)
        slide.add_element(
            TextElement(
                content=f"Short text {index}",
                bbox=[60.0, top, 560.0, top + 40.0],
                semantic_role="body",
            )
        )
    for index in range(4):
        left = 700.0 + (index % 2) * 250.0
        top = 120.0 + (index // 2) * 220.0
        slide.add_element(
            ImageElement(
                content="[IMAGE]",
                path=__file__,
                bbox=[left, top, left + 180.0, top + 160.0],
            )
        )
    presentation.slides.append(slide)

    summary = _presentation_summary(presentation, "req-fallback", "sample.pdf", "output-folder")

    assert summary["slides"][0]["default_render_mode"] == "image_fallback"


def test_generate_pptx_renders_multiline_title_as_multiple_paragraphs():
    presentation = Presentation()
    slide = Slide(page_id=1, width=1280.0, height=720.0)
    slide.add_element(
        TextElement(
            content="Big Title\nSecond Line",
            bbox=[100.0, 80.0, 700.0, 220.0],
            semantic_role="title",
            line_texts=["Big Title", "Second Line"],
            line_font_sizes=[28.0, 24.0],
        )
    )
    presentation.slides.append(slide)

    output_path = Path(generate_pptx(presentation, template_key="default", request_id="test-title-multiline"))
    ppt = PptxPresentation(str(output_path))
    text_shapes = [shape for shape in ppt.slides[0].shapes if hasattr(shape, "text") and shape.text]

    assert text_shapes
    assert text_shapes[0].text == "Big Title\nSecond Line"
    assert len(text_shapes[0].text_frame.paragraphs) == 2


def test_generate_pptx_keeps_multiline_body_as_single_paragraph_textbox():
    presentation = Presentation()
    slide = Slide(page_id=1, width=1280.0, height=720.0)
    slide.add_element(
        TextElement(
            content="Body first line\nBody second line",
            bbox=[100.0, 200.0, 900.0, 360.0],
            semantic_role="body",
            line_texts=["Body first line", "Body second line"],
            line_font_sizes=[18.0, 18.0],
        )
    )
    presentation.slides.append(slide)

    output_path = Path(generate_pptx(presentation, template_key="default", request_id="test-body-multiline"))
    ppt = PptxPresentation(str(output_path))
    text_shapes = [shape for shape in ppt.slides[0].shapes if hasattr(shape, "text") and shape.text]

    assert text_shapes
    assert text_shapes[0].text.replace("\v", "\n") == "Body first line\nBody second line"
    assert len(text_shapes[0].text_frame.paragraphs) == 1


def test_generate_pptx_uses_contain_transform_for_mismatched_aspect_ratio():
    presentation = Presentation(
        style_profile=DocumentStyleProfile(page_width=1200.0, page_height=900.0)
    )
    slide = Slide(page_id=1, width=1600.0, height=900.0)
    slide.add_element(
        TextElement(
            content="Mismatch",
            bbox=[160.0, 90.0, 800.0, 270.0],
            semantic_role="body",
        )
    )
    presentation.slides.append(slide)

    output_path = Path(generate_pptx(presentation, template_key="default", request_id="test-aspect-mismatch"))
    ppt = PptxPresentation(str(output_path))
    text_shapes = [shape for shape in ppt.slides[0].shapes if hasattr(shape, "text") and shape.text]

    assert text_shapes
    expected_left, expected_top, expected_width, expected_height = _element_geometry(
        slide.elements[0],
        slide.width,
        slide.height,
        ppt.slide_width,
        ppt.slide_height,
    )
    assert text_shapes[0].left == expected_left
    assert text_shapes[0].top == expected_top
    assert text_shapes[0].width == expected_width
    assert text_shapes[0].height == expected_height
    assert expected_top > 0


def test_generate_pptx_preserves_inline_text_styles():
    presentation = Presentation()
    slide = Slide(page_id=1, width=1280.0, height=720.0)
    slide.add_element(
        TextElement(
            content="Styled text",
            bbox=[100.0, 100.0, 600.0, 180.0],
            semantic_role="body",
            bold=True,
            italic=True,
            underline=True,
            strikethrough=True,
        )
    )
    presentation.slides.append(slide)

    output_path = Path(generate_pptx(presentation, template_key="default", request_id="test-inline-styles"))
    ppt = PptxPresentation(str(output_path))
    text_shapes = [shape for shape in ppt.slides[0].shapes if hasattr(shape, "text") and shape.text]

    assert text_shapes
    run = text_shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.italic is True
    assert run.font.underline is True
    assert 'strike="sngStrike"' in run._r.xml


def test_generate_pptx_separates_image_assets_from_text_when_available(tmp_path):
    image_path = tmp_path / "asset.png"
    Image.new("RGB", (160, 90), color=(24, 120, 200)).save(image_path)

    source_pdf_path = Path(__file__).resolve().parents[1] / "uploads" / "TheLastLeaf.pdf"
    assert source_pdf_path.exists(), f"Missing source PDF fixture: {source_pdf_path}"

    presentation = Presentation()
    slide = Slide(page_id=1, width=1280.0, height=720.0)
    slide.add_element(TextElement(content="Editable title", bbox=[80.0, 40.0, 520.0, 130.0], semantic_role="title"))
    slide.add_element(ImageElement(content="[IMAGE]", path=str(image_path), bbox=[700.0, 140.0, 1080.0, 420.0]))
    presentation.slides.append(slide)

    output_path = Path(
        generate_pptx(
            presentation,
            template_key="default",
            request_id="test-separate-image-assets",
            source_pdf_path=str(source_pdf_path),
        )
    )

    ppt = PptxPresentation(str(output_path))
    shapes = list(ppt.slides[0].shapes)
    picture_shapes = [shape for shape in shapes if shape.shape_type.name == "PICTURE"]
    text_shapes = [shape for shape in shapes if hasattr(shape, "text") and shape.text]

    assert text_shapes
    assert len(picture_shapes) == 1
    assert picture_shapes[0].width < ppt.slide_width
    assert picture_shapes[0].height < ppt.slide_height


def test_build_text_masked_picture_stream_masks_overlapping_text_regions(tmp_path):
    image_path = tmp_path / "masked-source.png"
    source_image = Image.new("RGBA", (200, 120), color=(232, 236, 240, 255))
    draw = ImageDraw.Draw(source_image)
    draw.rectangle((0, 0, 60, 40), fill=(18, 18, 18, 255))
    source_image.save(image_path)

    stream = _build_text_masked_picture_stream(
        str(image_path),
        [0.0, 0.0, 200.0, 120.0],
        [[0.0, 0.0, 60.0, 40.0]],
    )

    assert stream is not None
    masked_image = Image.open(stream).convert("RGBA")
    patched_pixel = masked_image.getpixel((10, 10))
    untouched_pixel = masked_image.getpixel((120, 60))

    assert patched_pixel[:3] != (255, 255, 255)
    assert abs(patched_pixel[0] - 232) <= 12
    assert abs(patched_pixel[1] - 236) <= 12
    assert abs(patched_pixel[2] - 240) <= 12
    assert untouched_pixel[:3] == (232, 236, 240)


def test_generate_pptx_applies_explicit_font_family_to_text_runs():
    presentation = Presentation()
    slide = Slide(page_id=1, width=1280.0, height=720.0)
    slide.add_element(
        TextElement(
            content="字体测试",
            bbox=[100.0, 100.0, 500.0, 180.0],
            semantic_role="title",
            font_name="Microsoft YaHei",
            color="#112233",
            bold=True,
        )
    )
    presentation.slides.append(slide)

    output_path = Path(generate_pptx(presentation, template_key="default", request_id="test-font-family"))
    ppt = PptxPresentation(str(output_path))
    text_shapes = [shape for shape in ppt.slides[0].shapes if hasattr(shape, "text") and shape.text]

    assert text_shapes
    run = text_shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Microsoft YaHei"
    assert text_shapes[0].text == "字体测试"

if __name__ == "__main__":
    test_generate_pptx_simple()
