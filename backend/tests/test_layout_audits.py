from __future__ import annotations

from pathlib import Path
from functools import lru_cache
from uuid import uuid4

from pptx import Presentation as PptxPresentation

from app.services.mineru_service import process_pdf
from app.services.parser_service import parse_mineru_output
from app.services.ppt_gen_service import generate_pptx

BASE_DIR = Path(__file__).resolve().parents[1]
OPENCLAW_PDF = BASE_DIR / "uploads" / "openclaw.pdf"
OVERLAP_THRESHOLD = 0.5
OVERFLOW_THRESHOLD_PT = 10.0


def _overlap_area(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> int:
    left = max(a[0], b[0])
    top = max(a[1], b[1])
    right = min(a[2], b[2])
    bottom = min(a[3], b[3])
    if right <= left or bottom <= top:
        return 0
    return (right - left) * (bottom - top)


def _shape_bbox(shape) -> tuple[int, int, int, int]:
    return (
        int(shape.left),
        int(shape.top),
        int(shape.left + shape.width),
        int(shape.top + shape.height),
    )


def _shape_text(shape) -> str:
    return " ".join(shape.text.replace("\v", "\n").split())


@lru_cache(maxsize=1)
def _run_openclaw_pipeline() -> Path:
    assert OPENCLAW_PDF.exists(), f"Missing fixture PDF: {OPENCLAW_PDF}"

    request_id = f"layout-audit-{uuid4().hex[:8]}"
    process_result = process_pdf(str(OPENCLAW_PDF), request_id=request_id)
    assert process_result["status"] == "success", process_result

    output_folder = Path(process_result["output_folder"])
    presentation = parse_mineru_output(str(output_folder))
    pptx_path = Path(
        generate_pptx(
            presentation,
            template_key="default",
            request_id=request_id,
            source_pdf_path=str(OPENCLAW_PDF),
        )
    )
    assert pptx_path.exists(), f"Generated PPTX missing: {pptx_path}"
    return pptx_path


def test_openclaw_text_boxes_have_no_severe_overlap():
    pptx_path = _run_openclaw_pipeline()
    prs = PptxPresentation(str(pptx_path))
    failures: list[str] = []

    for page_idx, slide in enumerate(list(prs.slides), start=1):
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_shapes.append((shape, _shape_text(shape), _shape_bbox(shape)))

        for a_idx in range(len(text_shapes)):
            for b_idx in range(a_idx + 1, len(text_shapes)):
                a_shape, a_text, a_bbox = text_shapes[a_idx]
                b_shape, b_text, b_bbox = text_shapes[b_idx]
                overlap = _overlap_area(a_bbox, b_bbox)
                if overlap <= 0:
                    continue
                area_a = max((a_bbox[2] - a_bbox[0]) * (a_bbox[3] - a_bbox[1]), 1)
                area_b = max((b_bbox[2] - b_bbox[0]) * (b_bbox[3] - b_bbox[1]), 1)
                ratio = overlap / min(area_a, area_b)
                if ratio >= OVERLAP_THRESHOLD:
                    failures.append(
                        f"slide {page_idx} ratio={ratio:.2f} A={a_text[:80]!r} B={b_text[:80]!r} Abox={a_bbox} Bbox={b_bbox}"
                    )

    assert not failures, "Severe text overlap detected in generated PPTX: " + "; ".join(failures[:10])


def test_openclaw_text_boxes_have_no_severe_overflow():
    pptx_path = _run_openclaw_pipeline()
    prs = PptxPresentation(str(pptx_path))
    failures: list[str] = []

    for page_idx, slide in enumerate(list(prs.slides), start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text:
                continue
            try:
                text_frame = shape.text_frame
            except Exception:
                continue
            paragraphs = list(text_frame.paragraphs)
            line_count = max(len(paragraphs), 1)
            font_sizes: list[float] = []
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    if run is not None and run.font is not None and run.font.size is not None:
                        font_sizes.append(float(run.font.size.pt))
            if not font_sizes:
                continue
            max_font_size = max(font_sizes)
            height_pt = float(shape.height) / 12700.0
            needed_pt = max_font_size * line_count * 1.18 + 2
            overflow_pt = needed_pt - height_pt
            if overflow_pt > OVERFLOW_THRESHOLD_PT:
                failures.append(
                    f"slide {page_idx} overflow={overflow_pt:.1f} height={height_pt:.1f} need={needed_pt:.1f} lines={line_count} size={max_font_size:.1f} text={_shape_text(shape)[:120]!r}"
                )

    assert not failures, "Severe text overflow detected in generated PPTX: " + "; ".join(failures[:10])
