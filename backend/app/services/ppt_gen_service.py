import io
import re
import statistics

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt, Inches
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image, ImageChops, ImageDraw, ImageFilter, ImageStat
from app.core.models import Presentation as DomainPresentation
from app.core.models import TextElement, DocumentStyleProfile
from app.core.config import settings
from pathlib import Path

try:
    import pypdfium2 as pdfium
except ImportError:  # pragma: no cover - optional runtime dependency
    pdfium = None

# Minimalist PPT Generator
BASE_DIR = Path(__file__).resolve().parents[2]
TEMPLATE_DIR = BASE_DIR / "assets" / "templates"
OUTPUT_DIR = BASE_DIR / "mineru_output"


def _hex_to_rgb(color: str | None, fallback: tuple[int, int, int]) -> RGBColor:
    if not color:
        return RGBColor(*fallback)
    normalized = color.lstrip("#")
    if len(normalized) != 6:
        return RGBColor(*fallback)
    try:
        return RGBColor(int(normalized[0:2], 16), int(normalized[2:4], 16), int(normalized[4:6], 16))
    except ValueError:
        return RGBColor(*fallback)


def _effective_style_profile(presentation: DomainPresentation) -> DocumentStyleProfile:
    if presentation.style_profile:
        return presentation.style_profile
    return DocumentStyleProfile()


def _build_base_presentation(presentation: DomainPresentation) -> Presentation:
    prs = Presentation()
    profile = _effective_style_profile(presentation)
    aspect_ratio = profile.page_width / profile.page_height if profile.page_height else (16 / 9)
    slide_height_inches = 7.5
    slide_width_inches = slide_height_inches * aspect_ratio
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    return prs


def _remove_existing_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _blank_layout(prs: Presentation):
    return min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))


def _remove_slide_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        element = shape.element
        element.getparent().remove(element)


def _is_cover_slide(d_slide) -> bool:
    if getattr(d_slide, "archetype", None) == "cover_split":
        return True
    text_elements = [element for element in d_slide.elements if element.type == "text" and getattr(element, "content", "").strip()]
    image_elements = [element for element in d_slide.elements if element.type in {"image", "table"}]
    return d_slide.page_id == 1 and len(text_elements) <= 3 and len(image_elements) == 1


def _alignment_value(value: str | None):
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _set_run_font_name(run, font_name: str) -> None:
    if not font_name:
        return
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        font = r_pr.find(qn(tag))
        if font is None:
            font = OxmlElement(tag)
            r_pr.append(font)
        font.set("typeface", font_name)


def _apply_run_inline_styles(run, elem: TextElement) -> None:
    italic = getattr(elem, "italic", None)
    underline = getattr(elem, "underline", None)
    strikethrough = getattr(elem, "strikethrough", None)

    if italic is not None:
        run.font.italic = bool(italic)
    if underline is not None:
        run.font.underline = bool(underline)
    if strikethrough is not None:
        r_pr = run._r.get_or_add_rPr()
        if bool(strikethrough):
            r_pr.set("strike", "sngStrike")
        else:
            r_pr.attrib.pop("strike", None)


def _infer_font_name(text: str, preferred: str | None = None) -> str:
    if preferred:
        return preferred
    if any("\u4e00" <= ch <= "\u9fff" for ch in text):
        return "Microsoft YaHei"
    return "Arial"


def _normalized_text(value: str | None) -> str:
    return " ".join((value or "").split())


def _box_height_points(textbox) -> float:
    return textbox.height / 12700 if textbox.height else 0.0


def _resolved_font_size_pt(textbox, elem: TextElement, fallback: float) -> float:
    explicit = getattr(elem, "font_size", None)
    if explicit:
        return max(float(explicit), 8.0)
    line_count = max(len(getattr(elem, "line_texts", None) or [line for line in (elem.content or "").splitlines() if line.strip()]), 1)
    box_height = _box_height_points(textbox)
    if box_height <= 0:
        return fallback
    estimated = box_height * 0.72 if line_count == 1 else box_height / (line_count * 1.35)
    return max(min(estimated, box_height * 0.85 if box_height else estimated), 8.0)


def _constrain_font_size_to_fit(textbox, elem: TextElement, base_size: float) -> float:
    """Shrink font size so that the longest line fits inside the textbox width
    and all lines fit inside the textbox height, preventing auto-wrap overflow."""
    # Leave ~15 % horizontal padding for margins / inter-character spacing
    width_pt = float(textbox.width) / 12700.0 * 0.85
    height_pt = float(textbox.height) / 12700.0

    line_texts = getattr(elem, "line_texts", None) or []
    if not line_texts:
        content = elem.content or ""
        if content:
            line_texts = [line for line in content.splitlines() if line.strip()]
    if not line_texts:
        return base_size

    role = getattr(elem, "semantic_role", None)
    if role == "title":
        spacing = settings.PPT_LINE_SPACING_TITLE
    elif role == "subtitle":
        spacing = settings.PPT_LINE_SPACING_SUBTITLE
    else:
        spacing = settings.PPT_LINE_SPACING_BODY

    num_lines = len(line_texts)

    # Height constraint: always apply so total lines fit in the box
    max_size_by_height = base_size
    if num_lines > 0 and height_pt > 0:
        max_size_by_height = height_pt / (num_lines * spacing)

    # Width constraint: for EVERY line in the block, ensure the text fits
    # inside the textbox width without PowerPoint auto-wrapping.
    # We take the most restrictive line so the unified block font size keeps
    # all original lines on single visual rows.
    max_size_by_width = base_size
    max_width_units = 0.0
    for line in line_texts:
        # Conservative estimates for PowerPoint rendering:
        # CJK full-width chars ~1.0× font pt; Latin ~0.55×
        units = sum(1.0 if "\u4e00" <= ch <= "\u9fff" else 0.55 for ch in line)
        max_width_units = max(max_width_units, units)
    if max_width_units > 0:
        max_size_by_width = width_pt / max_width_units

    constrained = min(base_size, max_size_by_width, max_size_by_height)
    return max(constrained, 6.0)


def _paragraph_font_size_pt(textbox, elem: TextElement, paragraph_index: int, fallback: float) -> float:
    line_font_sizes = getattr(elem, "line_font_sizes", None) or []
    if line_font_sizes:
        base_size = max(float(max(line_font_sizes)), 8.0)
    else:
        base_size = _resolved_font_size_pt(textbox, elem, fallback)

    return _constrain_font_size_to_fit(textbox, elem, base_size)


def _is_bullet_line(line: str) -> bool:
    """Return True if the line starts with a bullet/list prefix character."""
    return bool(re.search(r'^[\·\•\-\*\◦\‣\⁃\▪\▫\→\⇒\■\□]\s*', line.strip()))


def _strip_bullet_prefix(line: str) -> str:
    """Remove bullet prefix character from a line."""
    return re.sub(r'^[\·\•\-\*\◦\‣\⁃\▪\▫\→\⇒\■\□]\s*', '', line.strip())


def _set_paragraph_bullet(paragraph, char: str = "•") -> None:
    """Set a custom bullet character on a PowerPoint paragraph via OOXML."""
    pPr = paragraph._p.get_or_add_pPr()

    # Remove any existing bullet-none element
    bu_none = pPr.find(qn("a:buNone"))
    if bu_none is not None:
        pPr.remove(bu_none)

    # Remove existing buChar elements to avoid duplicates
    for existing in pPr.findall(qn("a:buChar")):
        pPr.remove(existing)

    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", char)
    pPr.append(bu_char)


def _should_word_wrap_text(elem: TextElement, content: str, line_texts: list[str]) -> bool:
    role = getattr(elem, "semantic_role", None)
    if role not in {"body", "subtitle", "caption"}:
        return False
    # If the original PDF already broke this text into multiple lines,
    # honour that by allowing word wrap.  If it was a single line, keep it
    # on one line and let the constrained font size handle the fit.
    if line_texts and len(line_texts) == 1 and "\n" not in content:
        return False
    if len(content) >= 18:
        return True
    return False


def _apply_text_style(textbox, elem: TextElement, is_cover_slide: bool) -> None:
    text_frame = textbox.text_frame
    role = getattr(elem, "semantic_role", None)
    content = (elem.content or "").strip()
    slide = textbox.part.slide
    profile = getattr(slide, "_pdf2ppt_style_profile", None)
    archetype = getattr(slide, "_pdf2ppt_archetype", "single_visual_explainer")
    profile_font_name = getattr(getattr(profile, "body_style", None), "font_name", None)
    resolved_font_name = _infer_font_name(content, getattr(elem, "font_name", None) or profile_font_name)
    is_short_year_label = archetype == "two_column_compare" and role == "subtitle" and content.isdigit() and len(content) <= 8

    title_color = _hex_to_rgb(getattr(elem, "color", None), (17, 17, 17))
    subtitle_color = _hex_to_rgb(getattr(elem, "color", None), (48, 114, 208) if is_short_year_label else (68, 68, 68))
    accent_color = _hex_to_rgb(getattr(elem, "color", None), (17, 17, 17))

    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        runs = [run for run in (paragraph.runs or []) if run.font]
        if not runs:
            runs = [paragraph.add_run()]
        paragraph.alignment = _alignment_value(getattr(elem, "align", None) or getattr(getattr(profile, "body_style", None), "align", "left"))
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

        if is_cover_slide and role == "title":
            target_size = _paragraph_font_size_pt(textbox, elem, paragraph_index, 30) * 0.9
            is_split_parenthetical_line = (paragraph_index > 0 and (paragraph.text or "").strip().startswith(("（", "("))) or content.startswith(("（", "("))
            for run in runs:
                run.font.size = Pt(target_size)
                run.font.bold = False if is_split_parenthetical_line else True
                run.font.color.rgb = subtitle_color if is_split_parenthetical_line else accent_color
                _set_run_font_name(run, _infer_font_name(content, getattr(elem, "font_name", None) or getattr(getattr(profile, "title_style", None), "font_name", None) or resolved_font_name))
                _apply_run_inline_styles(run, elem)
            paragraph.alignment = _alignment_value(getattr(elem, "align", None) or "left")
        elif is_cover_slide and role in {"subtitle", "body"} and len(content) <= 40:
            target_size = _paragraph_font_size_pt(textbox, elem, paragraph_index, 18) * 0.9
            for run in runs:
                run.font.size = Pt(target_size)
                run.font.bold = bool(getattr(elem, "bold", None))
                run.font.color.rgb = subtitle_color
                _set_run_font_name(run, _infer_font_name(content, getattr(elem, "font_name", None) or getattr(getattr(profile, "subtitle_style", None), "font_name", None) or resolved_font_name))
                _apply_run_inline_styles(run, elem)
            paragraph.alignment = _alignment_value(getattr(elem, "align", None) or "left")
        elif role == "title":
            target_size = _paragraph_font_size_pt(textbox, elem, paragraph_index, getattr(getattr(profile, "title_style", None), "font_size", 24) or 24)
            is_split_parenthetical_line = (paragraph_index > 0 and (paragraph.text or "").strip().startswith(("（", "("))) or content.startswith(("（", "("))
            for run in runs:
                run.font.size = Pt(target_size)
                run.font.bold = False if is_split_parenthetical_line else (True if getattr(elem, "bold", None) is None else bool(getattr(elem, "bold", None)))
                run.font.color.rgb = subtitle_color if is_split_parenthetical_line else title_color
                _set_run_font_name(run, _infer_font_name(content, getattr(elem, "font_name", None) or getattr(getattr(profile, "title_style", None), "font_name", None) or resolved_font_name))
                _apply_run_inline_styles(run, elem)
            paragraph.alignment = _alignment_value(getattr(elem, "align", None) or getattr(getattr(profile, "title_style", None), "align", "left"))
        elif role == "subtitle":
            target_size = _paragraph_font_size_pt(textbox, elem, paragraph_index, getattr(getattr(profile, "subtitle_style", None), "font_size", 18) or 18)
            for run in runs:
                run.font.size = Pt(target_size)
                run.font.bold = bool(getattr(elem, "bold", None))
                run.font.color.rgb = subtitle_color
                _set_run_font_name(run, _infer_font_name(content, getattr(elem, "font_name", None) or getattr(getattr(profile, "subtitle_style", None), "font_name", None) or resolved_font_name))
                _apply_run_inline_styles(run, elem)
            paragraph.alignment = _alignment_value(getattr(elem, "align", None) or getattr(getattr(profile, "subtitle_style", None), "align", "left"))
        elif role == "caption":
            target_size = _paragraph_font_size_pt(textbox, elem, paragraph_index, getattr(getattr(profile, "caption_style", None), "font_size", 12) or 12)
            for run in runs:
                run.font.size = Pt(target_size)
                run.font.bold = bool(getattr(elem, "bold", None))
                run.font.color.rgb = subtitle_color
                _set_run_font_name(run, _infer_font_name(content, getattr(elem, "font_name", None) or getattr(getattr(profile, "caption_style", None), "font_name", None) or resolved_font_name))
                _apply_run_inline_styles(run, elem)
            paragraph.alignment = _alignment_value(getattr(elem, "align", None) or getattr(getattr(profile, "caption_style", None), "align", "left"))
        else:
            target_size = _paragraph_font_size_pt(textbox, elem, paragraph_index, getattr(getattr(profile, "body_style", None), "font_size", 16) or 16)
            for run in runs:
                run.font.size = Pt(target_size)
                run.font.bold = bool(getattr(elem, "bold", None))
                run.font.color.rgb = title_color
                _set_run_font_name(run, resolved_font_name)
                _apply_run_inline_styles(run, elem)

        first_run = runs[0]
        first_run_size_pt = float(first_run.font.size.pt) if first_run.font.size is not None else None
        if archetype == "single_visual_explainer" and role == "title" and first_run_size_pt is not None:
            for run in runs:
                run.font.size = Pt(max(first_run_size_pt * 0.92, 10.0))
        if archetype == "single_visual_explainer" and role in {"subtitle", "body", "caption"} and first_run_size_pt is not None:
            for run in runs:
                run.font.size = Pt(max(first_run_size_pt * 0.93, 8.0))

        if archetype == "infographic_node_map" and role == "caption":
            for run in runs:
                run.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
        if archetype == "infographic_node_map" and role in {"subtitle", "body", "caption"} and first_run_size_pt is not None:
            for run in runs:
                run.font.size = Pt(max(first_run_size_pt * 0.85, 8.0))
        if archetype == "policy_text_heavy" and role == "title":
            for run in runs:
                run.font.size = Pt(26)

        unified_size_pt = float(first_run.font.size.pt) if first_run.font.size is not None else None
        if len(text_frame.paragraphs) > 1:
            if role == "title":
                paragraph.space_after = Pt(max(unified_size_pt * 0.22, 2.0)) if unified_size_pt else Pt(2.0)
            elif role in {"subtitle", "body", "caption"}:
                paragraph.space_after = Pt(max(unified_size_pt * 0.10, 1.0)) if unified_size_pt else Pt(1.0)

        font_size = first_run.font.size
        if font_size is not None:
            if role == "title":
                line_spacing_multiplier = settings.PPT_LINE_SPACING_TITLE
            elif role == "subtitle":
                line_spacing_multiplier = settings.PPT_LINE_SPACING_SUBTITLE
            else:
                line_spacing_multiplier = settings.PPT_LINE_SPACING_BODY
            paragraph.line_spacing = Pt(max(float(font_size.pt) * line_spacing_multiplier, 1.0))


def _apply_archetype_layout_hints(slide, d_slide, profile: DocumentStyleProfile) -> None:
    slide._pdf2ppt_style_profile = profile
    slide._pdf2ppt_archetype = getattr(d_slide, "archetype", "single_visual_explainer")


def _sorted_elements_for_render(d_slide):
    def _rank(element):
        bbox = element.bbox or [0, 0, 0, 0]
        order = 0 if element.type in {"image", "table"} else 1
        return (order, bbox[1], bbox[0])

    return sorted(d_slide.elements, key=_rank)


def _element_geometry(elem, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int) -> tuple[int, int, int, int]:
    x0, y0, x1, y1 = elem.bbox
    w = x1 - x0
    h = y1 - y0
    safe_pdf_w = max(float(pdf_w), 1.0)
    safe_pdf_h = max(float(pdf_h), 1.0)
    safe_ppt_w = max(int(ppt_w), 1)
    safe_ppt_h = max(int(ppt_h), 1)

    scale = min(safe_ppt_w / safe_pdf_w, safe_ppt_h / safe_pdf_h)
    scaled_page_w = safe_pdf_w * scale
    scaled_page_h = safe_pdf_h * scale
    offset_x = max((safe_ppt_w - scaled_page_w) / 2.0, 0.0)
    offset_y = max((safe_ppt_h - scaled_page_h) / 2.0, 0.0)

    left = int(round(offset_x + (x0 * scale)))
    top = int(round(offset_y + (y0 * scale)))
    width = max(int(round(w * scale)), 1)
    height = max(int(round(h * scale)), 1)
    return left, top, width, height


def _slide_page_index(d_slide) -> int:
    page_id = int(getattr(d_slide, "page_id", 1) or 1)
    return max(page_id - 1, 0)


def _render_text_element(slide, elem, left: int, top: int, width: int, height: int, is_cover_slide: bool) -> None:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.fill.background()
    textbox.line.fill.background()
    tf = textbox.text_frame
    role = getattr(elem, "semantic_role", None)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    archetype = getattr(slide, "_pdf2ppt_archetype", "single_visual_explainer")
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if role in {"body", "caption"} or (archetype == "infographic_node_map" and role == "subtitle") else MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    line_texts = [line for line in (getattr(elem, "line_texts", None) or []) if line.strip()]
    is_bullet_block = any(_is_bullet_line(line) for line in line_texts) if line_texts else _is_bullet_line(elem.content or "")

    if is_bullet_block:
        # Bullet item: merge continuation lines into one paragraph per bullet,
        # strip original bullet prefix, and add PowerPoint bullet character.
        first_paragraph = tf.paragraphs[0]
        bullet_texts: list[str] = []
        current_parts: list[str] = []

        for line in line_texts:
            if _is_bullet_line(line):
                if current_parts:
                    bullet_texts.append(" ".join(current_parts))
                current_parts = [_strip_bullet_prefix(line)]
            else:
                stripped = line.strip()
                if stripped:
                    current_parts.append(stripped)
        if current_parts:
            bullet_texts.append(" ".join(current_parts))

        if bullet_texts:
            first_paragraph.text = bullet_texts[0]
            _set_paragraph_bullet(first_paragraph, "•")
            for text in bullet_texts[1:]:
                paragraph = tf.add_paragraph()
                paragraph.text = text
                _set_paragraph_bullet(paragraph, "•")
    else:
        content_text = (elem.content or "").replace("\n", "\v")
        normalized_content = _normalized_text(content_text.replace("\v", "\n"))
        normalized_line_texts = _normalized_text("\n".join(line_texts)) if line_texts else ""
        should_render_lines_as_paragraphs = len(line_texts) > 1 and normalized_line_texts == normalized_content and role in {"title", "subtitle"}
        if should_render_lines_as_paragraphs:
            first_paragraph = tf.paragraphs[0]
            first_paragraph.text = line_texts[0]
            for line_text in line_texts[1:]:
                paragraph = tf.add_paragraph()
                paragraph.text = line_text
        else:
            tf.text = content_text
    tf.word_wrap = _should_word_wrap_text(elem, elem.content or "", line_texts)
    _apply_text_style(textbox, elem, is_cover_slide)


def _iter_text_overlay_units(elem: TextElement, archetype: str | None = None) -> list[TextElement]:
    return [elem]


def _render_text_elements_only(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool) -> None:
    archetype = getattr(d_slide, "archetype", None)
    for elem in _sorted_elements_for_render(d_slide):
        if elem.type != "text" or not elem.bbox:
            continue
        for overlay_elem in _iter_text_overlay_units(elem, archetype):
            if not overlay_elem.bbox:
                continue
            left, top, width, height = _element_geometry(overlay_elem, pdf_w, pdf_h, ppt_w, ppt_h)
            _render_text_element(slide, overlay_elem, left, top, width, height, is_cover_slide)


def _render_page_image(source_pdf_document, page_index: int, page_image_cache: dict[int, Image.Image]) -> Image.Image | None:
    cached = page_image_cache.get(page_index)
    if cached is not None:
        return cached

    if source_pdf_document is None or pdfium is None:
        return None

    try:
        page = source_pdf_document[page_index]
        bitmap = page.render(scale=8.0)
        pil_image = bitmap.to_pil()
    except Exception:
        return None

    page_image_cache[page_index] = pil_image
    return pil_image


def _mask_text_regions_in_image(image: Image.Image, image_bbox: list[float], text_bboxes: list[list[float]]) -> Image.Image | None:
    if len(image_bbox) != 4 or not text_bboxes:
        return None

    try:
        rgba = image.convert("RGBA")
        width, height = rgba.size
        if width <= 0 or height <= 0:
            return None

        image_left, image_top, image_right, image_bottom = [float(value) for value in image_bbox]
        image_width = max(image_right - image_left, 1.0)
        image_height = max(image_bottom - image_top, 1.0)
        has_overlap = False

        for text_bbox in text_bboxes:
            intersection = _bbox_intersection(image_bbox, text_bbox)
            if intersection is None:
                continue

            has_overlap = True
            overlap_left, overlap_top, overlap_right, overlap_bottom = intersection
            pixel_left = int(round(((overlap_left - image_left) / image_width) * width))
            pixel_top = int(round(((overlap_top - image_top) / image_height) * height))
            pixel_right = int(round(((overlap_right - image_left) / image_width) * width))
            pixel_bottom = int(round(((overlap_bottom - image_top) / image_height) * height))

            pixel_left = max(0, min(width - 1, pixel_left))
            pixel_top = max(0, min(height - 1, pixel_top))
            pixel_right = max(pixel_left + 1, min(width, pixel_right))
            pixel_bottom = max(pixel_top + 1, min(height, pixel_bottom))

            padding = 2
            pixel_left = max(0, pixel_left - padding)
            pixel_top = max(0, pixel_top - padding)
            pixel_right = min(width, pixel_right + padding)
            pixel_bottom = min(height, pixel_bottom + padding)
            fill_color = _sample_local_background_color(rgba, (pixel_left, pixel_top, pixel_right, pixel_bottom))
            if fill_color is None:
                continue
            rgba = _paint_background_patch(rgba, (pixel_left, pixel_top, pixel_right, pixel_bottom), fill_color)

        if not has_overlap:
            return None

        return rgba
    except Exception:
        return None

    page_image_cache[page_index] = pil_image
    return pil_image


def _render_picture_element(slide, elem, left: int, top: int, width: int, height: int, image_size_cache: dict[str, tuple[int, int]], source_pdf_document=None, page_index: int | None = None, pdf_w: float | None = None, pdf_h: float | None = None, page_image_cache: dict[int, Image.Image] | None = None) -> None:
    img_path = getattr(elem, "path", None)
    if img_path and Path(img_path).exists():
        text_bboxes = [text_element.bbox for text_element in getattr(slide, "_pdf2ppt_text_elements", []) if getattr(text_element, "bbox", None)]
        picture_source = _build_text_masked_picture_stream(img_path, list(getattr(elem, "bbox", []) or []), text_bboxes)
        _add_picture_cover(slide, img_path, left, top, width, height, image_size_cache, picture_source=picture_source)
        return

    if source_pdf_document is None or page_index is None or pdf_w is None or pdf_h is None:
        return

    if page_image_cache is None:
        page_image_cache = {}
    page_image = _render_page_image(source_pdf_document, page_index, page_image_cache)
    if page_image is None:
        return

    text_bboxes = [text_element.bbox for text_element in getattr(slide, "_pdf2ppt_text_elements", []) if getattr(text_element, "bbox", None)]

    bbox = getattr(elem, "bbox", None) or []
    if len(bbox) != 4:
        return

    page_width, page_height = page_image.size
    x0, y0, x1, y1 = [float(value) for value in bbox]
    crop_left = max(0, min(page_width, int(round((x0 / max(pdf_w, 1.0)) * page_width))))
    crop_top = max(0, min(page_height, int(round((y0 / max(pdf_h, 1.0)) * page_height))))
    crop_right = max(crop_left + 1, min(page_width, int(round((x1 / max(pdf_w, 1.0)) * page_width))))
    crop_bottom = max(crop_top + 1, min(page_height, int(round((y1 / max(pdf_h, 1.0)) * page_height))))

    try:
        cropped = page_image.crop((crop_left, crop_top, crop_right, crop_bottom))
    except Exception:
        return

    masked_crop = _mask_text_regions_in_image(cropped, [x0, y0, x1, y1], text_bboxes)
    if masked_crop is not None:
        cropped = masked_crop

    stream = io.BytesIO()
    cropped.save(stream, format="PNG")
    stream.seek(0)
    slide.shapes.add_picture(stream, left, top, width=width, height=height)


def _page_snapshot_candidate(d_slide):
    visual_elements = [element for element in d_slide.elements if element.type in {"image", "table"} and getattr(element, "path", None)]
    if not visual_elements:
        return None

    def _area(element):
        bbox = element.bbox or [0, 0, 0, 0]
        if len(bbox) != 4:
            return 0
        return max((bbox[2] - bbox[0]) * (bbox[3] - bbox[1]), 0)

    return max(visual_elements, key=_area)


def _slide_has_visual_blocks(d_slide) -> bool:
    return any(element.type in {"image", "table"} for element in d_slide.elements)


def _has_resolvable_visual_assets(d_slide) -> bool:
    for element in d_slide.elements:
        if element.type not in {"image", "table"}:
            continue
        img_path = getattr(element, "path", None)
        if img_path and Path(img_path).exists():
            return True
    return False


def _snapshot_area_ratio(d_slide) -> float:
    snapshot = _page_snapshot_candidate(d_slide)
    if snapshot is None or not snapshot.bbox or len(snapshot.bbox) != 4:
        return 0.0
    page_area = max((d_slide.width or 1280.0) * (d_slide.height or 720.0), 1.0)
    bbox = snapshot.bbox
    snapshot_area = max((bbox[2] - bbox[0]) * (bbox[3] - bbox[1]), 0.0)
    return snapshot_area / page_area


def _render_page_snapshot(slide, d_slide, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]]) -> bool:
    snapshot = _page_snapshot_candidate(d_slide)
    if snapshot is None:
        return False

    img_path = getattr(snapshot, "path", None)
    if not img_path or not Path(img_path).exists():
        return False

    _add_picture_cover(slide, img_path, 0, 0, ppt_w, ppt_h, image_size_cache)

    title_elements = [
        element
        for element in _sorted_elements_for_render(d_slide)
        if element.type == "text" and getattr(element, "semantic_role", None) in {"title", "subtitle"}
    ]

    for element in title_elements[:2]:
        if not element.bbox:
            continue
        left, top, width, height = _element_geometry(element, d_slide.width or 1280.0, d_slide.height or 720.0, ppt_w, ppt_h)
        _render_text_element(slide, element, left, top, width, height, is_cover_slide)

    return True


def _render_pdf_page_snapshot(slide, source_pdf_document, page_index: int, ppt_w: int, ppt_h: int) -> bool:
    if source_pdf_document is None or pdfium is None:
        return False

    try:
        page = source_pdf_document[page_index]
        bitmap = page.render(scale=8.0)
        pil_image = bitmap.to_pil()
    except Exception:
        return False

    stream = io.BytesIO()
    pil_image.save(stream, format="PNG")
    stream.seek(0)
    slide.shapes.add_picture(stream, 0, 0, width=ppt_w, height=ppt_h)
    return True


def _is_ppt_like_slide(d_slide, render_mode: str, source_pdf_document) -> bool:
    return False


def _render_title_overlays_only(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool) -> None:
    title_elements = [
        element
        for element in _sorted_elements_for_render(d_slide)
        if element.type == "text" and getattr(element, "semantic_role", None) in {"title", "subtitle"}
    ]
    for element in title_elements[:3]:
        for overlay_elem in _iter_text_overlay_units(element):
            if not overlay_elem.bbox:
                continue
            left, top, width, height = _element_geometry(overlay_elem, pdf_w, pdf_h, ppt_w, ppt_h)
            _render_text_element(slide, overlay_elem, left, top, width, height, is_cover_slide)


def _render_ppt_like_fidelity(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]], source_pdf_document=None) -> bool:
    render_mode = getattr(d_slide, "render_mode", "auto")
    if not _is_ppt_like_slide(d_slide, render_mode, source_pdf_document):
        return False
    if not _render_page_snapshot(slide, d_slide, ppt_w, ppt_h, is_cover_slide, image_size_cache):
        return False
    _render_text_elements_only(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide)
    return True


def _render_generic_archetype(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]], source_pdf_document=None) -> None:
    page_index = _slide_page_index(d_slide)
    page_image_cache: dict[int, Image.Image] = {}
    slide._pdf2ppt_text_elements = [element for element in d_slide.elements if element.type == "text"]
    for elem in _sorted_elements_for_render(d_slide):
        if not elem.bbox:
            continue
        if elem.type == "text":
            for overlay_elem in _iter_text_overlay_units(elem):
                if not overlay_elem.bbox:
                    continue
                left, top, width, height = _element_geometry(overlay_elem, pdf_w, pdf_h, ppt_w, ppt_h)
                _render_text_element(slide, overlay_elem, left, top, width, height, is_cover_slide)
        elif elem.type in {"image", "table"}:
            left, top, width, height = _element_geometry(elem, pdf_w, pdf_h, ppt_w, ppt_h)
            _render_picture_element(
                slide,
                elem,
                left,
                top,
                width,
                height,
                image_size_cache,
                None if source_pdf_document is None else source_pdf_document,
                page_index,
                pdf_w,
                pdf_h,
                page_image_cache,
            )


def _render_cover_split(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, image_size_cache: dict[str, tuple[int, int]]) -> None:
    _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, True, image_size_cache, None)


def _render_roadmap_overview(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]]) -> None:
    _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache, None)


def _render_two_column_compare(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]]) -> None:
    _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache, None)


def _render_policy_text_heavy(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]]) -> None:
    _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache, None)


def _render_closing_statement(slide, d_slide, pdf_w: float, pdf_h: float, ppt_w: int, ppt_h: int, is_cover_slide: bool, image_size_cache: dict[str, tuple[int, int]]) -> None:
    _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache, None)


def _render_slide_by_archetype(
    slide,
    d_slide,
    pdf_w: float,
    pdf_h: float,
    ppt_w: int,
    ppt_h: int,
    is_cover_slide: bool,
    render_mode: str,
    image_size_cache: dict[str, tuple[int, int]],
    source_pdf_document=None,
) -> None:
    page_index = _slide_page_index(d_slide)
    text_elements = [element for element in d_slide.elements if element.type == "text" and getattr(element, "content", "").strip()]
    visual_elements = [element for element in d_slide.elements if element.type in {"image", "table"}]

    if render_mode in {"image_fallback", "hybrid_overlay"} and source_pdf_document is not None:
        if _render_pdf_page_snapshot(slide, source_pdf_document, page_index, ppt_w, ppt_h):
            return

    if render_mode == "editable":
        _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache, source_pdf_document)
        return

    if _render_ppt_like_fidelity(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache):
        return

    _render_generic_archetype(slide, d_slide, pdf_w, pdf_h, ppt_w, ppt_h, is_cover_slide, image_size_cache, source_pdf_document)


def _get_image_size(img_path: str, image_size_cache: dict[str, tuple[int, int]]) -> tuple[int, int]:
    cached = image_size_cache.get(img_path)
    if cached:
        return cached

    with Image.open(img_path) as image:
        size = image.size

    image_size_cache[img_path] = size
    return size


def _compute_content_crop(img_path: str) -> tuple[float, float, float, float]:
    """Return (left, top, right, bottom) crop ratios to trim background/whitespace.

    Uses the alpha channel when available; otherwise samples border pixels to
    estimate the background colour and crops to the bounding box of pixels
    that differ from it by more than a threshold.
    """
    try:
        with Image.open(img_path) as img:
            rgba = img.convert("RGBA")
            w, h = rgba.size
            if w <= 0 or h <= 0:
                return (0.0, 0.0, 0.0, 0.0)

            r, g, b, a = rgba.split()

            # 1. Alpha mask (fast LUT)
            alpha_mask = a.point(lambda v: 255 if v > 10 else 0)

            # 2. Background colour from unique border samples
            samples = set()
            for x in range(w):
                samples.add(rgba.getpixel((x, 0)))
                samples.add(rgba.getpixel((x, h - 1)))
            for y in range(h):
                samples.add(rgba.getpixel((0, y)))
                samples.add(rgba.getpixel((w - 1, y)))
            if not samples:
                return (0.0, 0.0, 0.0, 0.0)

            bg_r = int(statistics.median(p[0] for p in samples))
            bg_g = int(statistics.median(p[1] for p in samples))
            bg_b = int(statistics.median(p[2] for p in samples))

            # 3. Difference from background per channel (fast LUT)
            threshold = 20
            r_diff = r.point(lambda v: abs(v - bg_r))
            g_diff = g.point(lambda v: abs(v - bg_g))
            b_diff = b.point(lambda v: abs(v - bg_b))

            # 4. Combine: max(diff_r, diff_g, diff_b)
            diff = ImageChops.lighter(ImageChops.lighter(r_diff, g_diff), b_diff)
            colour_mask = diff.point(lambda v: 255 if v > threshold else 0)

            # 5. Final mask = colour_mask OR alpha_mask
            final_mask = ImageChops.lighter(colour_mask, alpha_mask)

            bbox = final_mask.getbbox()
            if not bbox:
                return (0.0, 0.0, 0.0, 0.0)

            # Sanity check: ignore if detected content is < 2 % of total area
            content_area = (bbox[2] - bbox[0]) * (bbox[3] - bbox[1])
            if content_area < 0.02 * w * h:
                return (0.0, 0.0, 0.0, 0.0)

            return (
                bbox[0] / w,
                bbox[1] / h,
                (w - bbox[2]) / w,
                (h - bbox[3]) / h,
            )
    except Exception:
        return (0.0, 0.0, 0.0, 0.0)


def _bbox_intersection(a: list[float] | tuple[float, float, float, float], b: list[float] | tuple[float, float, float, float]) -> tuple[float, float, float, float] | None:
    if len(a) != 4 or len(b) != 4:
        return None

    left = max(float(a[0]), float(b[0]))
    top = max(float(a[1]), float(b[1]))
    right = min(float(a[2]), float(b[2]))
    bottom = min(float(a[3]), float(b[3]))
    if right <= left or bottom <= top:
        return None
    return left, top, right, bottom


def _sample_local_background_color(rgba: Image.Image, fill_bbox: tuple[int, int, int, int], border_size: int = 8) -> tuple[int, int, int] | None:
    width, height = rgba.size
    left, top, right, bottom = fill_bbox
    regions: list[Image.Image] = []

    if top > 0:
        regions.append(rgba.crop((left, max(0, top - border_size), right, top)))
    if bottom < height:
        regions.append(rgba.crop((left, bottom, right, min(height, bottom + border_size))))
    if left > 0:
        regions.append(rgba.crop((max(0, left - border_size), top, left, bottom)))
    if right < width:
        regions.append(rgba.crop((right, top, min(width, right + border_size), bottom)))

    weighted_channels = [0.0, 0.0, 0.0]
    total_weight = 0.0
    for region in regions:
        if region.width <= 0 or region.height <= 0:
            continue
        stat = ImageStat.Stat(region)
        weight = float(region.width * region.height)
        for index in range(3):
            weighted_channels[index] += float(stat.mean[index]) * weight
        total_weight += weight

    if total_weight > 0:
        return tuple(int(round(weighted_channels[index] / total_weight)) for index in range(3))

    samples: list[tuple[int, int, int, int]] = []
    for x in range(width):
        samples.append(rgba.getpixel((x, 0)))
        samples.append(rgba.getpixel((x, height - 1)))
    for y in range(height):
        samples.append(rgba.getpixel((0, y)))
        samples.append(rgba.getpixel((width - 1, y)))

    if not samples:
        return None

    return (
        int(round(statistics.median(pixel[0] for pixel in samples))),
        int(round(statistics.median(pixel[1] for pixel in samples))),
        int(round(statistics.median(pixel[2] for pixel in samples))),
    )


def _paint_background_patch(rgba: Image.Image, fill_bbox: tuple[int, int, int, int], fill_color: tuple[int, int, int], feather_radius: float = 2.0) -> Image.Image:
    patch = Image.new("RGBA", rgba.size, fill_color + (255,))
    mask = Image.new("L", rgba.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rectangle(fill_bbox, fill=255)
    if feather_radius > 0:
        mask = mask.filter(ImageFilter.GaussianBlur(radius=feather_radius))
    return Image.composite(patch, rgba, mask)


def _build_text_masked_picture_stream(img_path: str, image_bbox: list[float], text_bboxes: list[list[float]]) -> io.BytesIO | None:
    if len(image_bbox) != 4 or not text_bboxes:
        return None

    try:
        with Image.open(img_path) as img:
            rgba = img.convert("RGBA")
            width, height = rgba.size
            if width <= 0 or height <= 0:
                return None

            image_left, image_top, image_right, image_bottom = [float(value) for value in image_bbox]
            image_width = max(image_right - image_left, 1.0)
            image_height = max(image_bottom - image_top, 1.0)
            has_overlap = False

            for text_bbox in text_bboxes:
                intersection = _bbox_intersection(image_bbox, text_bbox)
                if intersection is None:
                    continue

                has_overlap = True
                overlap_left, overlap_top, overlap_right, overlap_bottom = intersection
                pixel_left = int(round(((overlap_left - image_left) / image_width) * width))
                pixel_top = int(round(((overlap_top - image_top) / image_height) * height))
                pixel_right = int(round(((overlap_right - image_left) / image_width) * width))
                pixel_bottom = int(round(((overlap_bottom - image_top) / image_height) * height))

                pixel_left = max(0, min(width - 1, pixel_left))
                pixel_top = max(0, min(height - 1, pixel_top))
                pixel_right = max(pixel_left + 1, min(width, pixel_right))
                pixel_bottom = max(pixel_top + 1, min(height, pixel_bottom))

                padding = 2
                pixel_left = max(0, pixel_left - padding)
                pixel_top = max(0, pixel_top - padding)
                pixel_right = min(width, pixel_right + padding)
                pixel_bottom = min(height, pixel_bottom + padding)
                fill_color = _sample_local_background_color(rgba, (pixel_left, pixel_top, pixel_right, pixel_bottom))
                if fill_color is None:
                    continue
                rgba = _paint_background_patch(rgba, (pixel_left, pixel_top, pixel_right, pixel_bottom), fill_color)

            if not has_overlap:
                return None

            stream = io.BytesIO()
            rgba.save(stream, format="PNG")
            stream.seek(0)
            return stream
    except Exception:
        return None


def _add_picture_cover(slide, img_path: str, left: int, top: int, width: int, height: int, image_size_cache: dict[str, tuple[int, int]], picture_source=None) -> None:
    image_width, image_height = _get_image_size(img_path, image_size_cache)
    source = picture_source if picture_source is not None else img_path

    if image_width <= 0 or image_height <= 0:
        slide.shapes.add_picture(source, left, top, width=width, height=height)
        return

    picture = slide.shapes.add_picture(source, left, top, width=width, height=height)

    # Content-aware crop: trim whitespace/background so the visible image
    # occupies less of the rectangular MinerU bbox, reducing overlap with
    # neighbouring text elements.
    c_left, c_top, c_right, c_bottom = _compute_content_crop(img_path)

    content_width = image_width * (1 - c_left - c_right)
    content_height = image_height * (1 - c_top - c_bottom)
    frame_aspect = width / height if height else 1
    content_aspect = content_width / content_height if content_height else frame_aspect

    if content_aspect > frame_aspect:
        aspect_crop = (1 - (frame_aspect / content_aspect)) / 2
        total_crop_left = c_left + (1 - c_left - c_right) * aspect_crop
        total_crop_right = c_right + (1 - c_left - c_right) * aspect_crop
        total_crop_top = c_top
        total_crop_bottom = c_bottom
    else:
        aspect_crop = (1 - (content_aspect / frame_aspect)) / 2 if frame_aspect else 0
        total_crop_top = c_top + (1 - c_top - c_bottom) * aspect_crop
        total_crop_bottom = c_bottom + (1 - c_top - c_bottom) * aspect_crop
        total_crop_left = c_left
        total_crop_right = c_right

    picture.crop_left = total_crop_left
    picture.crop_right = total_crop_right
    picture.crop_top = total_crop_top
    picture.crop_bottom = total_crop_bottom


def generate_pptx(
    presentation: DomainPresentation,
    template_key: str = "default",
    request_id: str | None = None,
    render_mode_overrides: dict[int, str] | None = None,
    source_pdf_path: str | None = None,
) -> str:
    """
    Convert Domain Presentation to PPTX file.
    Returns path to the generated file.
    """
    overrides = render_mode_overrides or {}
    image_size_cache: dict[str, tuple[int, int]] = {}
    source_pdf_document = None
    if source_pdf_path and pdfium is not None:
        source_path = Path(source_pdf_path)
        if source_path.exists():
            source_pdf_document = pdfium.PdfDocument(str(source_path))

    prs = _build_base_presentation(presentation)
    profile = _effective_style_profile(presentation)

    # Layout mapping (Naive)
    title_layout = _blank_layout(prs)
    content_layout = _blank_layout(prs)  # Title and Content

    # Iterate domain slides
    for i, d_slide in enumerate(presentation.slides):
        # Create PPT slide
        # If it's the first page, use Title layout? Or just content for all?
        # Let's use Content layout for now.
        layout = content_layout if i > 0 else title_layout
        slide = prs.slides.add_slide(layout)
        _remove_slide_placeholders(slide)
        is_cover_slide = _is_cover_slide(d_slide)
        _apply_archetype_layout_hints(slide, d_slide, profile)

        # Mapping Elements
        # Mineru elements: Text, Image, Table

        # Title handling: Mineru might have extracted a Title element.
        # But we don't have explicit "Title" type in Element yet.
        # We rely on TextElement content or just generic placement.

        # For structure-preserving, we might want to use Blank layout and place elements
        # at exact coordinates.
        # But Mineru bbox is [x0, y0, x1, y1] typically in PDF points (72dpi?).
        # PPTX uses Inches/EMUs.
        # Scale Factor: PDF Width / PPT Width

        # Use the actual template slide size instead of a hardcoded assumption.
        ppt_w = prs.slide_width
        ppt_h = prs.slide_height

        # Determine PDF scale
        # Domain Slide has width/height?
        pdf_w = d_slide.width if d_slide.width else 595.0  # A4 point width approx
        pdf_h = d_slide.height if d_slide.height else 842.0

        # We'll use a safer approach: Place elements relative to slide size using %
        # bbox is usually [x0, y0, x1, y1]

        render_mode = overrides.get(d_slide.page_id) or getattr(d_slide, "render_mode", None) or "auto"
        _render_slide_by_archetype(
            slide,
            d_slide,
            pdf_w,
            pdf_h,
            ppt_w,
            ppt_h,
            is_cover_slide,
            render_mode,
            image_size_cache,
            source_pdf_document,
        )

    # Save output
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_filename = f"{request_id or 'generated'}_presentation.pptx"
    output_path = OUTPUT_DIR / output_filename
    prs.save(str(output_path))

    
    return str(output_path)
