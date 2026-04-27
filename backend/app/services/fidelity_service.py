from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from app.core.models import Presentation, TextElement


@dataclass
class FidelityMatch:
    parsed_index: int
    generated_index: int
    text_score: float
    layout_score: float
    font_size_score: float
    font_family_score: float
    color_score: float
    paragraph_score: float
    alignment_score: float
    role: str | None = None
    expected_font_size: float | None = None
    actual_font_size: float | None = None
    expected_alignment: str | None = None
    actual_alignment: str | None = None


@dataclass
class SlideFidelityReport:
    page_id: int
    score: float
    text_score: float
    layout_score: float
    font_size_score: float
    font_family_score: float
    color_score: float
    image_score: float
    paragraph_score: float
    alignment_score: float
    hierarchy_score: float
    matched_text_count: int
    total_text_count: int
    matched_picture_count: int
    total_picture_count: int


@dataclass
class FidelityReport:
    score: float
    text_score: float
    layout_score: float
    font_size_score: float
    font_family_score: float
    color_score: float
    image_score: float
    paragraph_score: float
    alignment_score: float
    hierarchy_score: float
    slide_count_score: float
    full_page_snapshot_penalty: float
    matched_text_count: int
    total_text_count: int
    matched_picture_count: int
    total_picture_count: int
    slide_reports: list[SlideFidelityReport] = field(default_factory=list)
    matches: list[FidelityMatch] = field(default_factory=list)


def _normalized_text(value: str) -> str:
    return " ".join((value or "").split())


def _shape_text(shape) -> str:
    return getattr(shape, "text", "") if hasattr(shape, "text") else ""


def _shape_bbox(shape, slide_width: int, slide_height: int) -> tuple[float, float, float, float]:
    return (
        shape.left / max(slide_width, 1),
        shape.top / max(slide_height, 1),
        shape.width / max(slide_width, 1),
        shape.height / max(slide_height, 1),
    )


def _element_bbox(element: TextElement, slide_width: float, slide_height: float) -> tuple[float, float, float, float]:
    x0, y0, x1, y1 = element.bbox or [0, 0, 0, 0]
    return (
        x0 / max(slide_width, 1.0),
        y0 / max(slide_height, 1.0),
        (x1 - x0) / max(slide_width, 1.0),
        (y1 - y0) / max(slide_height, 1.0),
    )


def _bbox_distance(actual: tuple[float, float, float, float], expected: tuple[float, float, float, float]) -> float:
    return sum(abs(a - b) for a, b in zip(actual, expected))


def _score_from_distance(distance: float, tolerance: float = 0.08) -> float:
    if distance <= tolerance:
        return 1.0
    return max(0.0, 1.0 - ((distance - tolerance) / max(tolerance, 1e-6)))


def _score_font_size(actual: float | None, expected: float | None) -> float:
    if not actual or not expected:
        return 0.7
    ratio = actual / max(expected, 1e-6)
    if 0.9 <= ratio <= 1.1:
        return 1.0
    if 0.8 <= ratio <= 1.2:
        return 0.85
    return max(0.0, 1.0 - abs(1.0 - ratio))


def _score_font_family(actual: str | None, expected: str | None, text: str) -> float:
    if expected:
        if actual and actual.lower() == expected.lower():
            return 1.0
        return 0.0
    if any("\u4e00" <= ch <= "\u9fff" for ch in text):
        return 1.0 if actual and actual.lower() == "microsoft yahei" else 0.5
    return 1.0 if actual and actual.lower() in {"arial", "aptos", "calibri"} else 0.5


def _score_color(actual: str | None, expected: str | None) -> float:
    if not expected:
        return 0.75
    if not actual:
        return 0.0
    return 1.0 if actual.lower() == expected.lower() else 0.5


def _score_alignment(actual: str | None, expected: str | None) -> float:
    if not expected or not actual:
        return 0.75
    if actual.lower() == expected.lower():
        return 1.0
    return 0.25


def _score_ratio(actual: float | None, expected: float | None, tolerance: float = 0.25) -> float:
    if not actual or not expected:
        return 0.75
    ratio = actual / max(expected, 1e-6)
    if 1.0 - tolerance <= ratio <= 1.0 + tolerance:
        return 1.0
    return max(0.0, 1.0 - abs(1.0 - ratio) / max(1.0 + tolerance, 1e-6))


def _median(values: list[float]) -> float | None:
    if not values:
        return None
    values = sorted(values)
    middle = len(values) // 2
    if len(values) % 2:
        return values[middle]
    return (values[middle - 1] + values[middle]) / 2.0


def _safe_run_font_size(shape) -> float | None:
    try:
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            font_size = paragraph.runs[0].font.size
            if font_size is not None:
                return float(font_size.pt)
    except Exception:
        return None


def _safe_paragraph_alignment(shape) -> str | None:
    try:
        paragraph = shape.text_frame.paragraphs[0]
        alignment = paragraph.alignment
        if alignment is None:
            return None
        if alignment == PP_ALIGN.LEFT:
            return "left"
        if alignment == PP_ALIGN.CENTER:
            return "center"
        if alignment == PP_ALIGN.RIGHT:
            return "right"
        if alignment == PP_ALIGN.JUSTIFY:
            return "justify"
    except Exception:
        return None
    return None


def _safe_run_font_name(shape) -> str | None:
    try:
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            return paragraph.runs[0].font.name
    except Exception:
        return None
    return None


def _safe_run_color(shape) -> str | None:
    try:
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            color = paragraph.runs[0].font.color
            if color and color.rgb is not None:
                return f"#{color.rgb}".lower()
    except Exception:
        return None
    return None


def _count_picture_shapes(slide) -> int:
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def _has_full_page_snapshot(slide, slide_width: int, slide_height: int) -> bool:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.left == 0 and shape.top == 0 and shape.width == slide_width and shape.height == slide_height:
            return True
    return False


def _match_text_elements(parsed_slide, generated_slide, slide_width: int, slide_height: int) -> tuple[list[FidelityMatch], int, int]:
    parsed_text_elements = [element for element in parsed_slide.elements if isinstance(element, TextElement) and element.bbox and (element.content or "").strip()]
    generated_text_shapes = [shape for shape in generated_slide.shapes if hasattr(shape, "text") and shape.text]
    unused_indexes: set[int] = set()
    matches: list[FidelityMatch] = []

    for parsed_index, element in enumerate(parsed_text_elements):
        expected_text = _normalized_text(element.content)
        expected_bbox = _element_bbox(element, parsed_slide.width, parsed_slide.height)
        candidates: list[tuple[int, Any, float]] = []
        for generated_index, shape in enumerate(generated_text_shapes):
            if generated_index in unused_indexes:
                continue
            if _normalized_text(_shape_text(shape)) != expected_text:
                continue
            actual_bbox = _shape_bbox(shape, slide_width, slide_height)
            candidates.append((generated_index, shape, _bbox_distance(actual_bbox, expected_bbox)))

        if not candidates:
            continue

        candidates.sort(key=lambda item: item[2])
        generated_index, shape, _ = candidates[0]
        unused_indexes.add(generated_index)

        actual_bbox = _shape_bbox(shape, slide_width, slide_height)
        text_score = 1.0
        layout_score = _score_from_distance(_bbox_distance(actual_bbox, expected_bbox))
        actual_font_size = _safe_run_font_size(shape)
        expected_font_size = element.font_size
        font_size_score = _score_font_size(actual_font_size, expected_font_size)
        font_family_score = _score_font_family(_safe_run_font_name(shape), getattr(element, "font_name", None), element.content)
        color_score = _score_color(_safe_run_color(shape), getattr(element, "color", None))
        paragraph_count = len(getattr(shape.text_frame, "paragraphs", [])) if hasattr(shape, "text_frame") else 1
        expected_lines = len(getattr(element, "line_texts", None) or [line for line in (element.content or "").splitlines() if line.strip()])
        paragraph_score = 1.0 if paragraph_count == expected_lines else max(0.0, 1.0 - abs(paragraph_count - expected_lines) * 0.25)
        actual_alignment = _safe_paragraph_alignment(shape)
        expected_alignment = getattr(element, "align", None)
        alignment_score = _score_alignment(actual_alignment, expected_alignment)

        matches.append(
            FidelityMatch(
                parsed_index=parsed_index,
                generated_index=generated_index,
                text_score=text_score,
                layout_score=layout_score,
                font_size_score=font_size_score,
                font_family_score=font_family_score,
                color_score=color_score,
                paragraph_score=paragraph_score,
                alignment_score=alignment_score,
                role=getattr(element, "semantic_role", None),
                expected_font_size=expected_font_size,
                actual_font_size=actual_font_size,
                expected_alignment=expected_alignment,
                actual_alignment=actual_alignment,
            )
        )

    return matches, len(parsed_text_elements), len(generated_text_shapes)


def _score_slide_components(
    text_score: float,
    layout_score: float,
    font_size_score: float,
    font_family_score: float,
    color_score: float,
    image_score: float,
    paragraph_score: float,
    alignment_score: float,
    hierarchy_score: float,
) -> float:
    return (
        text_score * 0.17
        + layout_score * 0.27
        + font_size_score * 0.17
        + font_family_score * 0.11
        + color_score * 0.08
        + image_score * 0.10
        + paragraph_score * 0.05
        + alignment_score * 0.05
        + hierarchy_score * 0.03
    ) * 100.0


def _safe_slide_collection(ppt: PptxPresentation) -> list[Any]:
    return list(ppt.slides)


def evaluate_presentation_fidelity(parsed: Presentation, ppt: PptxPresentation | str | Path) -> FidelityReport:
    if isinstance(ppt, (str, Path)):
        ppt = PptxPresentation(str(ppt))

    def _avg(values: Iterable[float], fallback: float = 1.0) -> float:
        values = list(values)
        return sum(values) / len(values) if values else fallback

    slides = _safe_slide_collection(ppt)
    parsed_slides = list(parsed.slides)
    paired_count = min(len(parsed_slides), len(slides))

    text_scores: list[float] = []
    layout_scores: list[float] = []
    font_size_scores: list[float] = []
    font_family_scores: list[float] = []
    color_scores: list[float] = []
    paragraph_scores: list[float] = []
    alignment_scores: list[float] = []
    hierarchy_scores: list[float] = []
    image_scores: list[float] = []
    matches: list[FidelityMatch] = []
    slide_reports: list[SlideFidelityReport] = []
    matched_text_count = 0
    total_text_count = 0
    matched_picture_count = 0
    total_picture_count = 0
    slide_count_score = 1.0 if len(slides) >= len(parsed_slides) else max(0.0, len(slides) / max(len(parsed_slides), 1))
    full_page_snapshot_penalty = 0.0

    for idx in range(paired_count):
        parsed_slide = parsed_slides[idx]
        generated_slide = slides[idx]
        slide_width = int(ppt.slide_width)
        slide_height = int(ppt.slide_height)

        slide_matches, slide_text_total, _ = _match_text_elements(parsed_slide, generated_slide, slide_width, slide_height)
        matches.extend(slide_matches)
        matched_text_count += len(slide_matches)
        total_text_count += slide_text_total

        if slide_matches:
            text_scores.extend([m.text_score for m in slide_matches])
            layout_scores.extend([m.layout_score for m in slide_matches])
            font_size_scores.extend([m.font_size_score for m in slide_matches])
            font_family_scores.extend([m.font_family_score for m in slide_matches])
            color_scores.extend([m.color_score for m in slide_matches])
            paragraph_scores.extend([m.paragraph_score for m in slide_matches])
            alignment_scores.extend([m.alignment_score for m in slide_matches])

            role_sizes: dict[str, list[float]] = {"title": [], "subtitle": [], "body": [], "caption": []}
            for match in slide_matches:
                if match.role and match.actual_font_size:
                    role_sizes.setdefault(match.role, []).append(match.actual_font_size)

            title_size = _median(role_sizes.get("title", []))
            subtitle_size = _median(role_sizes.get("subtitle", []))
            body_size = _median(role_sizes.get("body", []))
            caption_size = _median(role_sizes.get("caption", []))

            slide_role_scores: list[float] = []
            if title_size and body_size:
                slide_role_scores.append(_score_ratio(title_size / max(body_size, 1e-6), 1.5, 0.30))
            if subtitle_size and body_size:
                slide_role_scores.append(_score_ratio(subtitle_size / max(body_size, 1e-6), 1.15, 0.35))
            if body_size and caption_size:
                slide_role_scores.append(_score_ratio(body_size / max(caption_size, 1e-6), 1.2, 0.35))
            if slide_role_scores:
                hierarchy_scores.append(sum(slide_role_scores) / len(slide_role_scores))

        parsed_picture_count = sum(1 for element in parsed_slide.elements if getattr(element, "type", None) in {"image", "table"} and element.bbox)
        generated_picture_count = _count_picture_shapes(generated_slide)
        total_picture_count += parsed_picture_count
        matched_picture_count += min(parsed_picture_count, generated_picture_count)
        if parsed_picture_count:
            image_scores.append(min(1.0, generated_picture_count / parsed_picture_count))

        slide_text_score = _avg([m.text_score for m in slide_matches], 0.0) if slide_matches else 0.0
        slide_layout_score = _avg([m.layout_score for m in slide_matches], 0.0) if slide_matches else 0.0
        slide_font_size_score = _avg([m.font_size_score for m in slide_matches], 0.0) if slide_matches else 0.0
        slide_font_family_score = _avg([m.font_family_score for m in slide_matches], 0.0) if slide_matches else 0.0
        slide_color_score = _avg([m.color_score for m in slide_matches], 0.0) if slide_matches else 0.0
        slide_paragraph_score = _avg([m.paragraph_score for m in slide_matches], 1.0) if slide_matches else 1.0
        slide_alignment_score = _avg([m.alignment_score for m in slide_matches], 0.0) if slide_matches else 0.0
        slide_hierarchy_score = hierarchy_scores[-1] if hierarchy_scores else 1.0
        slide_image_score = min(1.0, generated_picture_count / parsed_picture_count) if parsed_picture_count else 1.0

        if _has_full_page_snapshot(generated_slide, slide_width, slide_height):
            full_page_snapshot_penalty += 0.10

        slide_score = max(
            0.0,
            min(
                100.0,
                _score_slide_components(
                    slide_text_score,
                    slide_layout_score,
                    slide_font_size_score,
                    slide_font_family_score,
                    slide_color_score,
                    slide_image_score,
                    slide_paragraph_score,
                    slide_alignment_score,
                    slide_hierarchy_score,
                ),
            ),
        )
        slide_reports.append(
            SlideFidelityReport(
                page_id=getattr(parsed_slide, "page_id", idx + 1),
                score=slide_score,
                text_score=slide_text_score,
                layout_score=slide_layout_score,
                font_size_score=slide_font_size_score,
                font_family_score=slide_font_family_score,
                color_score=slide_color_score,
                image_score=slide_image_score,
                paragraph_score=slide_paragraph_score,
                alignment_score=slide_alignment_score,
                hierarchy_score=slide_hierarchy_score,
                matched_text_count=len(slide_matches),
                total_text_count=slide_text_total,
                matched_picture_count=min(parsed_picture_count, generated_picture_count),
                total_picture_count=parsed_picture_count,
            )
        )

    text_score = _avg(text_scores, 0.0)
    layout_score = _avg(layout_scores, 0.0)
    font_size_score = _avg(font_size_scores, 0.0)
    font_family_score = _avg(font_family_scores, 0.0)
    color_score = _avg(color_scores, 0.0)
    image_score = _avg(image_scores, 1.0)
    paragraph_score = _avg(paragraph_scores, 1.0)
    alignment_score = _avg(alignment_scores, 0.0)
    hierarchy_score = _avg(hierarchy_scores, 1.0)

    weighted_score = (
        text_score * 0.16
        + layout_score * 0.25
        + font_size_score * 0.16
        + font_family_score * 0.10
        + color_score * 0.07
        + image_score * 0.09
        + paragraph_score * 0.05
        + alignment_score * 0.06
        + hierarchy_score * 0.04
        + slide_count_score * 0.02
    ) * 100.0

    score = max(0.0, min(100.0, weighted_score - (full_page_snapshot_penalty * 100.0)))

    return FidelityReport(
        score=score,
        text_score=text_score,
        layout_score=layout_score,
        font_size_score=font_size_score,
        font_family_score=font_family_score,
        color_score=color_score,
        image_score=image_score,
        paragraph_score=paragraph_score,
        alignment_score=alignment_score,
        hierarchy_score=hierarchy_score,
        slide_count_score=slide_count_score,
        full_page_snapshot_penalty=full_page_snapshot_penalty,
        matched_text_count=matched_text_count,
        total_text_count=total_text_count,
        matched_picture_count=matched_picture_count,
        total_picture_count=total_picture_count,
        slide_reports=slide_reports,
        matches=matches,
    )
