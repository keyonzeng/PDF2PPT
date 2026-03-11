from pptx import Presentation
import os

def create_simple_template(filename, title):
    prs = Presentation()
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Subtitle"
    
    # Save
    path = os.path.join("assets", "templates", filename)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"Created template: {path}")

if __name__ == "__main__":
    create_simple_template("default.pptx", "Default Template")
    create_simple_template("academic.pptx", "Academic Template")
    create_simple_template("minimal.pptx", "Minimal Template")
